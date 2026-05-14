"""
TARS v10.0.0-rc.1 — Presentation deck builder.

Cinematic dark theme, Interstellar-monolith inspired.
20 slides, all with speaker notes (~100 words each), ready for live demo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Palette ----------
BG          = RGBColor(0x07, 0x08, 0x0d)
PANEL       = RGBColor(0x0d, 0x11, 0x1b)
PANEL_2     = RGBColor(0x13, 0x18, 0x26)
LINE        = RGBColor(0x1f, 0x29, 0x3b)
WHITE       = RGBColor(0xff, 0xff, 0xff)
BODY        = RGBColor(0xcb, 0xd5, 0xe1)
MUTED       = RGBColor(0x94, 0xa3, 0xb8)
DIM         = RGBColor(0x64, 0x74, 0x8b)
INDIGO      = RGBColor(0x63, 0x66, 0xf1)
VIOLET      = RGBColor(0x8b, 0x5c, 0xf6)
CYAN        = RGBColor(0x06, 0xb6, 0xd4)
LIME        = RGBColor(0x84, 0xcc, 0x16)
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)
ROSE        = RGBColor(0xf4, 0x3f, 0x5e)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def color_line(shape, color, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        shp.adjustments[0] = 0.10
    if fill is None:
        shp.fill.background()
    else:
        set_fill(shp, fill)
    if line is None:
        no_line(shp)
    else:
        color_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, color=WHITE, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for (text, size, color, bold, italic, font) in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font or "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def background(slide):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    bg.shadow.inherit = False
    return bg


def vignette_corners(slide):
    r = Inches(0.06)
    add_rect(slide, Inches(0.45), Inches(0.45), r, r, fill=CYAN, rounded=True)
    add_rect(slide, SLIDE_W - Inches(0.51), Inches(0.45), r, r, fill=VIOLET, rounded=True)


def footer(slide, num, total=20):
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(6), Inches(0.3),
             "TARS  ·  v10.0.0-rc.1  ·  meeet.world",
             size=9, color=DIM, font="Consolas")
    add_text(slide, SLIDE_W - Inches(2.0), SLIDE_H - Inches(0.45),
             Inches(1.4), Inches(0.3),
             f"{num:02d} / {total:02d}",
             size=9, color=DIM, font="Consolas", align=PP_ALIGN.RIGHT)


def title_bar(slide, eyebrow, title, *, title_color=WHITE,
              title_size=36, italic_title=False):
    add_text(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.32),
             eyebrow.upper(), size=11, color=CYAN, bold=True,
             font="Consolas")
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12.2), Inches(0.9),
             title, size=title_size, color=title_color, bold=True,
             italic=italic_title, font="Calibri")


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ---------- Visual motif ----------

def draw_monolith(slide, cx, cy, h, w_ratio=0.32):
    w = int(h * w_ratio)
    x = int(cx - w / 2)
    y = int(cy - h / 2)
    body = add_rect(slide, x, y, w, h, fill=RGBColor(0x05, 0x05, 0x09))
    color_line(body, RGBColor(0x1a, 0x1f, 0x2e), 1.0)
    strip_w = int(w * 0.12)
    strip_x = int(x + w * 0.44)
    add_rect(slide, strip_x, y + int(h * 0.08), strip_w, int(h * 0.42), fill=CYAN)
    add_rect(slide, strip_x, y + int(h * 0.52), strip_w, int(h * 0.40), fill=VIOLET)
    add_rect(slide, x - Inches(0.05), y + int(h * 0.5) - Inches(0.04),
             w + Inches(0.10), Inches(0.08), fill=INDIGO)


def draw_ring(slide, cx, cy, radius, *, color=VIOLET, thickness_pt=1.5):
    x = cx - radius
    y = cy - radius
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, radius * 2, radius * 2)
    shp.fill.background()
    color_line(shp, color, thickness_pt)
    shp.shadow.inherit = False
    return shp


def hex_token(slide, cx, cy, size, label, sublabel, color):
    half = size / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, cx - half, cy - half, size, size)
    set_fill(shp, PANEL_2)
    color_line(shp, color, 1.5)
    shp.shadow.inherit = False
    add_text(slide, cx - half, cy - Inches(0.22), size, Inches(0.32),
             label, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             font="Calibri")
    add_text(slide, cx - half, cy + Inches(0.05), size, Inches(0.25),
             sublabel, size=8, color=color, align=PP_ALIGN.CENTER,
             font="Consolas")


def quad_card(slide, x, y, w, h, eyebrow, title, body, accent):
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, rounded=True)
    add_rect(slide, x + Inches(0.25), y + Inches(0.25),
             Inches(0.10), Inches(0.10), fill=accent, rounded=True)
    add_text(slide, x + Inches(0.45), y + Inches(0.22),
             w - Inches(0.6), Inches(0.3),
             eyebrow.upper(), size=10, color=accent, bold=True,
             font="Consolas")
    add_text(slide, x + Inches(0.30), y + Inches(0.65),
             w - Inches(0.6), Inches(0.5),
             title, size=20, color=WHITE, bold=True, font="Calibri")
    add_text(slide, x + Inches(0.30), y + Inches(1.20),
             w - Inches(0.6), h - Inches(1.3),
             body, size=13, color=BODY, font="Calibri", line_spacing=1.35)


def bullet_list(slide, x, y, w, h, items, *, size=15, color=BODY,
                bullet_color=CYAN, gap=Inches(0.42)):
    cy = y
    for txt in items:
        add_rect(slide, x, cy + Inches(0.10), Inches(0.10), Inches(0.10),
                 fill=bullet_color, rounded=True)
        add_text(slide, x + Inches(0.25), cy, w - Inches(0.3), gap,
                 txt, size=size, color=color, font="Calibri", line_spacing=1.3)
        cy += gap


def stat_callout(slide, x, y, w, h, number, label, color):
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, rounded=True)
    # Reserve top 65% for number, bottom 30% for label, no overlap
    num_h = h * 0.62
    add_text(slide, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2),
             num_h, number,
             size=46, color=color, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.1), y + num_h + Inches(0.05),
             w - Inches(0.2), h - num_h - Inches(0.1), label.upper(),
             size=9, color=MUTED, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# =====================================================================
# Build the deck
# =====================================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def slide_1():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)

    draw_monolith(s, SLIDE_W - Inches(2.6), Inches(3.75),
                  Inches(5.2), w_ratio=0.28)
    draw_ring(s, SLIDE_W - Inches(2.6), Inches(3.75), Inches(2.0),
              color=RGBColor(0x33, 0x36, 0x6f), thickness_pt=0.75)
    draw_ring(s, SLIDE_W - Inches(2.6), Inches(3.75), Inches(2.5),
              color=RGBColor(0x22, 0x24, 0x4c), thickness_pt=0.5)

    add_text(s, Inches(0.8), Inches(0.8), Inches(4), Inches(0.32),
             "TARS  ·  v10.0.0-rc.1", size=11, color=CYAN, bold=True,
             font="Consolas")

    add_text(s, Inches(0.6), Inches(1.4), Inches(8), Inches(2.6),
             "TARS", size=140, color=WHITE, bold=True, font="Calibri")

    add_runs(s, Inches(0.8), Inches(4.05), Inches(8.4), Inches(0.5),
             [("AI cockpit ", 22, WHITE, True, False, "Calibri"),
              ("for everything not code.", 22, VIOLET, True, True, "Calibri")])

    add_text(s, Inches(0.8), Inches(4.75), Inches(8.4), Inches(1.5),
             "Local-first. Voice-native. Receipt-anchored.\n"
             "Billed through meeet.world.",
             size=17, color=BODY, font="Calibri", line_spacing=1.40)

    add_rect(s, Inches(0.8), Inches(6.10),
             Inches(3.7), Inches(0.50),
             fill=PANEL, line=CYAN, rounded=True)
    add_text(s, Inches(0.8), Inches(6.13), Inches(3.7), Inches(0.4),
             "  v10.0.0-rc.1  ·  ships tomorrow",
             size=12, color=CYAN, bold=True, font="Consolas",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 1)
    set_notes(s,
        "30-second pitch. Hold for a beat. 'This is TARS. Cursor solved coding. "
        "We're doing it for everything else — wealth, health, business, creative "
        "work. One install on your Mac. Voice or chat. Local-first, so your data "
        "never has to leave the box. Every consequential action emits a "
        "cryptographic receipt anchored on Solana. Billed through meeet.world, "
        "my brother's domain. v10.0 release candidate is on the laptop in front "
        "of you. We ship general availability tomorrow. The next twenty minutes "
        "show you why this is a different shape of product than ChatGPT, "
        "Cursor, or Claude Desktop.'")
    return s


def slide_2():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "The problem",
              "Cursor solved coding.",
              title_size=42)
    add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.7),
             "The rest of your life is still ChatGPT tabs.",
             size=34, color=VIOLET, bold=True, italic=True, font="Calibri")

    x = Inches(0.6); w = Inches(3.95); y = Inches(2.85); h = Inches(3.65); gap = Inches(0.17)

    quad_card(s, x, y, w, h, "Friction",
              "7-12 AI tabs.\nZero memory.",
              "ChatGPT in one tab. Claude in another. Gemini, Notion AI, "
              "Linear AI, Granola, the calendar copilot. None of them "
              "remember context across tools. None of them take action — "
              "they're all chat surfaces.",
              CYAN)
    quad_card(s, x + (w + gap), y, w, h, "Lock-in",
              "Cursor is for devs.\nNothing else exists.",
              "Cursor proved the cockpit shape works — for engineers, "
              "inside an IDE. The other 80% of knowledge work has no "
              "equivalent: investing, ops, legal, health, creative, "
              "civic — all stuck in chat.",
              VIOLET)
    quad_card(s, x + 2 * (w + gap), y, w, h, "No receipts",
              "'Did the agent really do that?'",
              "Today the answer is 'check the screenshot.' No tampering "
              "detection, no third-party verification, no audit trail "
              "regulators or partners can trust. Agents act; nothing proves it.",
              ROSE)

    footer(s, 2)
    set_notes(s,
        "Open with the emotional hook. 'How many AI tabs do you have open "
        "right now? Be honest.' Pause. 'Mine is twelve. None of them talk "
        "to each other. None of them remember what I told the other one "
        "five minutes ago. None of them actually do things — they all just "
        "tell me how I could do things. Cursor fixed this. For engineers. "
        "Inside an IDE. The other eighty percent of my day — running my "
        "fund, managing my health, the family calendar — is still "
        "ChatGPT-tab-soup. And when an agent finally does take an action, "
        "I have no way to prove it actually happened.'")
    return s


def slide_3():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "The thesis",
              "What if you had a Cursor for everything?",
              title_size=36)

    add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
             "Investing.  Health.  Business.  Creative.  Civic.  Family.  Research.",
             size=18, color=BODY, italic=True, font="Calibri")

    cx = SLIDE_W / 2
    cy = Inches(4.85)
    size = Inches(1.55)
    dx = Inches(1.55)

    packs = [
        ("WEALTH",       "balance/anchor",  AMBER),
        ("HEALTH",       "vitals/coach",     LIME),
        ("BUSINESS",     "kpi/brief",        CYAN),
        ("ENTREPRENEUR", "outreach/pipe",    VIOLET),
        ("RESEARCH",     "arxiv/notes",      INDIGO),
        ("CIVIC",        "vote/lookup",      RGBColor(0x10, 0xb9, 0x81)),
        ("CREATIVE",     "brand/draft",      ROSE),
    ]

    top_y = cy - Inches(0.7)
    bot_y = cy + Inches(0.65)
    top = packs[:4]; bot = packs[4:]
    top_start = cx - (1.5 * dx)
    for i, (lbl, sub, col) in enumerate(top):
        hex_token(s, top_start + i * dx, top_y, size, lbl, sub, col)
    bot_start = cx - (1.0 * dx)
    for i, (lbl, sub, col) in enumerate(bot):
        hex_token(s, bot_start + i * dx, bot_y, size, lbl, sub, col)

    footer(s, 3)
    set_notes(s,
        "Vision slide. Speak slowly. 'What if you had a Cursor for "
        "investing? For your health? For your business? For creative work, "
        "civic engagement, family logistics, research? Same cockpit shape, "
        "different vocabulary, all on your machine. That's TARS. Seven "
        "domain packs shipped today: wealth, health, business, "
        "entrepreneur, research, civic — civic is free for every tier "
        "because public-records access shouldn't be paywalled — and "
        "creative. Each pack knows the vocabulary, the obvious next move, "
        "the connectors that matter. Not ask-me-anything. Domain-aware.'")
    return s


def slide_4():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "What TARS is", "Four pillars. One cockpit.",
              title_size=36)

    px = Inches(0.6); py = Inches(1.85); pw = Inches(6.05); ph = Inches(2.55); gap = Inches(0.18)

    quad_card(s, px, py, pw, ph,
              "Pillar 1", "Voice-first",
              "Wake-word → STT → router → action → TTS → receipt. "
              "Cmd+Shift+Space anywhere on macOS takes you from idle to "
              "voice command in 200ms. Cinematic monolith UI — not a "
              "chat window with a mic strapped on.",
              CYAN)
    quad_card(s, px + pw + gap, py, pw, ph,
              "Pillar 2", "Local-first",
              "Backend is a Python sidecar on 127.0.0.1:8765. Memory, "
              "receipts, vector index on disk. Keys in macOS Keychain. "
              "Air-gap the laptop and TARS still works for everything "
              "that doesn't strictly need a cloud LLM call.",
              VIOLET)
    quad_card(s, px, py + ph + gap, pw, ph,
              "Pillar 3", "Receipt-anchored",
              "Every consequential action emits a signed receipt. "
              "Receipts hash-chain. Daily Merkle root anchors on Solana "
              "as a memo. Public verifier replays the proof with no auth "
              "and no access to our database.",
              INDIGO)
    quad_card(s, px + pw + gap, py + ph + gap, pw, ph,
              "Pillar 4", "meeet.world-native",
              "Identity, billing, balance, entitlements, marketplace "
              "payouts, $MEEET economy, compliance telemetry — every "
              "cross-machine action flows through my brother's domain. "
              "TARS is the cockpit; meeet.world is the ledger and the bank.",
              ROSE)

    footer(s, 4)
    set_notes(s,
        "Four pillars. Read them in order, the order matters. "
        "'Voice-first because the cockpit is ambient — you don't open it, "
        "you summon it. Local-first because regulated industries can't "
        "send data to OpenAI by default and most of us would prefer not "
        "to either. Receipt-anchored because action-taking agents must be "
        "auditable; this is the missing primitive. And meeet.world-native "
        "because identity and money need a substrate, and my brother runs "
        "exactly that substrate. Four pillars, one cockpit. Every feature "
        "we ship has to pass all four checks.'")
    return s


def slide_5():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)

    add_text(s, Inches(0.6), Inches(0.55), Inches(10), Inches(0.32),
             "LIVE DEMO  ·  PRESENTER CUE", size=11, color=AMBER, bold=True,
             font="Consolas")

    add_text(s, Inches(0.6), Inches(1.05), Inches(12.2), Inches(2.4),
             "Now let me show you.", size=58, color=WHITE, bold=True,
             italic=True, font="Calibri")

    add_rect(s, Inches(0.6), Inches(3.65), Inches(5.2), Inches(2.3),
             fill=PANEL, line=CYAN, line_w=2.0, rounded=True)
    add_text(s, Inches(0.6), Inches(3.85), Inches(5.2), Inches(0.5),
             "DEMO", size=72, color=CYAN, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.05), Inches(5.2), Inches(0.5),
             "60 SECONDS TO FIRST AGENT ACTION",
             size=13, color=BODY, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(6.4), Inches(3.7), Inches(6.4), Inches(0.4),
             "BEFORE YOU CLICK", size=12, color=AMBER, bold=True,
             font="Consolas")
    bullet_list(s, Inches(6.4), Inches(4.1), Inches(6.4), Inches(2.6),
                ["TARS.app already open, sidecar :8765 green",
                 "TARS_DEMO_SEED=1 verified — fixtures loaded",
                 "Network OK · LLM_PROVIDER=openrouter live",
                 "Audience can see the monolith, the cyan strip is breathing"],
                size=13, gap=Inches(0.45))

    footer(s, 5)
    set_notes(s,
        "Stop. Take a breath. Make eye contact. 'OK — slides are useful, "
        "but you came to see whether the thing works. Let me show you.' "
        "Switch to the TARS.app window. Make sure the sidecar dot is "
        "green. Make sure the monolith is breathing — that one cyan-violet "
        "light strip slowly pulsing. If anything looks frozen, hit "
        "Cmd+Shift+Space to refocus. The next slide is the runbook in "
        "case something dies. If everything is healthy, hit the global "
        "shortcut and start the voice command. The clock starts now.")
    return s


def slide_6():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Demo script", "Five beats.  Sixty seconds.  One wow.",
              title_size=32)

    steps = [
        ("01", "LAUNCH",   "Cmd+Shift+Space",
         "Monolith fades in.\nCyan strip charges.", "~3s", CYAN),
        ("02", "VOICE",    "\"Compose a thank-you\nto last week's investors.\"",
         "STT waveform pulses.\nRouter picks Composer.", "~8s", VIOLET),
        ("03", "DIFF",     "Composer panel opens.\nMulti-file diff preview.",
         "Domain-pack aware.\nEntrepreneur vocab.", "~15s", INDIGO),
        ("04", "APPROVE",  "Click  Accept hunks.\nGmail draft created.",
         "Real OAuth.\nReal draft in Drafts folder.", "~10s", LIME),
        ("05", "AUDIT",    "Open Audit Explorer.\nReceipt #1284 is live.",
         "Hash chained, ready to\nMerkle-anchor on Solana.", "~5s", AMBER),
    ]

    col_w = Inches(2.45); col_h = Inches(4.55)
    start_x = Inches(0.5); y = Inches(1.95); gap = Inches(0.10)

    for i, (num, label, do, wow, t, color) in enumerate(steps):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, y, col_w, col_h, fill=PANEL, line=LINE, rounded=True)
        add_rect(s, x + Inches(0.2), y + Inches(0.2),
                 Inches(0.55), Inches(0.40), fill=color, rounded=True)
        add_text(s, x + Inches(0.2), y + Inches(0.22),
                 Inches(0.55), Inches(0.40),
                 num, size=14, color=BG, bold=True, font="Consolas",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y + Inches(0.22),
                 Inches(1.5), Inches(0.4),
                 label, size=15, color=WHITE, bold=True, font="Calibri",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.9),
                 col_w - Inches(0.4), Inches(0.3),
                 "ACTION", size=9, color=color, bold=True,
                 font="Consolas")
        add_text(s, x + Inches(0.2), y + Inches(1.20),
                 col_w - Inches(0.4), Inches(1.2),
                 do, size=12, color=WHITE, bold=True, font="Calibri",
                 line_spacing=1.3)
        add_text(s, x + Inches(0.2), y + Inches(2.5),
                 col_w - Inches(0.4), Inches(0.3),
                 "WOW MOMENT", size=9, color=color, bold=True,
                 font="Consolas")
        add_text(s, x + Inches(0.2), y + Inches(2.80),
                 col_w - Inches(0.4), Inches(1.3),
                 wow, size=11, color=BODY, italic=True, font="Calibri",
                 line_spacing=1.3)
        add_rect(s, x + Inches(0.2), y + col_h - Inches(0.55),
                 Inches(0.85), Inches(0.32), fill=PANEL_2, line=color,
                 rounded=True)
        add_text(s, x + Inches(0.2), y + col_h - Inches(0.55),
                 Inches(0.85), Inches(0.32),
                 t, size=10, color=color, bold=True, font="Consolas",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 6)
    set_notes(s,
        "Verbal script for the demo. Beat 1: 'Cmd+Shift+Space — and you "
        "can see the monolith come alive.' Pause for the breath. Beat 2: "
        "say the command verbatim — 'Compose a thank-you to last week's "
        "investors.' Don't improvise the phrase, the STT is trained on "
        "this one. Beat 3: 'See how the Composer knows the Entrepreneur "
        "pack vocabulary — the diff is already written for that audience.' "
        "Beat 4: hit Accept, point at Gmail in another window: 'real "
        "draft, real OAuth, no mocks.' Beat 5: open Audit. 'Receipt is "
        "live and Merkle-chained. Tomorrow it anchors on Solana.'")
    return s


def slide_7():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Architecture", "Local sidecar.  Receipt anchor.  Token rails.",
              title_size=30)

    y = Inches(2.3)
    nh = Inches(1.85)
    nw = Inches(2.2)

    nodes = [
        ("TARS.app",      "Tauri 2 desktop\nVoice cockpit\nCmd+Shift+Space",       CYAN,   Inches(0.6)),
        ("Backend",       "FastAPI :8765\n50 routers\nSQLite + sqlite-vec",        VIOLET, Inches(3.2)),
        ("meeet.world",   "Identity / Billing\n$MEEET economy\nMarketplace 70/30", INDIGO, Inches(5.8)),
        ("Solana",        "Memo program\nDaily Merkle root\nPublic verifier",      AMBER,  Inches(8.4)),
        ("LLM providers", "Anthropic / OpenAI\nOpenRouter / Ollama\nBYO key",      ROSE,   Inches(11.0)),
    ]
    for (title, body, color, x) in nodes:
        add_rect(s, x, y, nw, nh, fill=PANEL, line=color, line_w=1.5, rounded=True)
        add_rect(s, x + Inches(0.2), y + Inches(0.2),
                 Inches(0.12), Inches(0.12), fill=color, rounded=True)
        add_text(s, x + Inches(0.4), y + Inches(0.16),
                 nw - Inches(0.5), Inches(0.4),
                 title, size=14, color=WHITE, bold=True, font="Calibri")
        add_text(s, x + Inches(0.2), y + Inches(0.6),
                 nw - Inches(0.4), nh - Inches(0.7),
                 body, size=10, color=BODY, font="Consolas",
                 line_spacing=1.3)

    for i in range(4):
        x_from = Inches(0.6) + nw + i * Inches(2.6)
        chev = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  x_from, y + nh / 2 - Inches(0.10),
                                  Inches(0.40), Inches(0.20))
        set_fill(chev, DIM)
        no_line(chev)
        chev.shadow.inherit = False

    add_text(s, Inches(0.6), Inches(5.05), Inches(12.2), Inches(0.4),
             "DATA PATH", size=10, color=CYAN, bold=True, font="Consolas")
    add_text(s, Inches(0.6), Inches(5.40), Inches(12.2), Inches(0.5),
             "Voice/chat → router → agent → action → "
             "signed receipt → daily Merkle root → Solana memo.",
             size=14, color=BODY, font="Calibri")
    add_text(s, Inches(0.6), Inches(6.05), Inches(12.2), Inches(0.4),
             "AUTHORITY", size=10, color=VIOLET, bold=True, font="Consolas")
    add_text(s, Inches(0.6), Inches(6.40), Inches(12.2), Inches(0.5),
             "meeet.world holds identity, balance, entitlements, "
             "compliance telemetry — the only cloud surface.",
             size=14, color=BODY, font="Calibri")

    footer(s, 7)
    set_notes(s,
        "Technical credibility slide. Walk left to right. 'TARS.app is a "
        "Tauri 2 desktop wrapper — Rust shell, web frontend. It talks "
        "exclusively to a Python FastAPI sidecar on localhost 8765 — fifty "
        "routers, SQLite plus sqlite-vec for the vector index, everything "
        "on disk. The sidecar talks to meeet.world for identity, billing, "
        "and the $MEEET economy — that's the one cloud surface. From "
        "meeet.world we batch daily Merkle roots to Solana as memos. LLM "
        "providers are pluggable — Anthropic, OpenAI, OpenRouter, local "
        "Ollama, or bring your own key. No vendor lock-in by design.'")
    return s


def slide_8():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Scorecard", "Cursor parity + TARS-only edge.",
              title_size=32)

    stat_callout(s, Inches(9.4), Inches(1.55), Inches(1.75), Inches(1.4),
                 "15", "Cursor parity", LIME)
    stat_callout(s, Inches(11.25), Inches(1.55), Inches(1.75), Inches(1.4),
                 "10", "TARS-only", VIOLET)

    rows = [
        ("Code editor",                  "Y", "—", "by design"),
        ("Inline Tab completion",        "Y", "—", "wraps Cursor"),
        ("Composer (multi-file diff)",   "Y", "Y",      "W253"),
        ("Codebase index",               "Y", "Y",      "sqlite-vec"),
        ("@-mention context",            "Y", "Y",      "W240"),
        ("MCP servers panel",            "Y", "Y",      "W238"),
        ("Model switcher + cost label",  "Y", "Y",      "W237"),
        ("Usage console",                "Y", "Y",      "W235"),
        ("Magic-link auth",              "Y", "Y",      "W219"),
        ("Privacy mode",                 "Y", "Y",      "local-first"),
        ("Voice-first cockpit",          "—", "Y", "W220/W230"),
        ("Hash-chained receipts",        "—", "Y", "W67/W95"),
        ("Solana anchor",                "—", "Y", "W89"),
        ("7 domain packs",               "—", "Y", "wealth/health/etc"),
        ("$MEEET economy + marketplace", "—", "Y", "70/30 split"),
    ]

    tx = Inches(0.6); ty = Inches(3.05); tw = Inches(12.2)
    row_h = Inches(0.245); header_h = Inches(0.32)

    add_rect(s, tx, ty, tw, header_h, fill=PANEL_2, line=LINE)
    add_text(s, tx + Inches(0.2), ty, Inches(6), header_h,
             "CAPABILITY", size=10, color=CYAN, bold=True, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + Inches(6.3), ty, Inches(1.5), header_h,
             "CURSOR", size=10, color=CYAN, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + Inches(7.9), ty, Inches(1.5), header_h,
             "TARS", size=10, color=CYAN, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + Inches(9.6), ty, Inches(2.5), header_h,
             "NOTE", size=10, color=CYAN, bold=True, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)

    for i, (cap, cu, ta, note) in enumerate(rows):
        ry = ty + header_h + i * row_h
        bg_color = BG if i % 2 == 0 else PANEL
        add_rect(s, tx, ry, tw, row_h, fill=bg_color)
        add_text(s, tx + Inches(0.2), ry, Inches(6), row_h,
                 cap, size=11, color=WHITE, font="Calibri",
                 anchor=MSO_ANCHOR.MIDDLE)
        cu_color = LIME if cu == "Y" else DIM
        ta_color = LIME if ta == "Y" else DIM
        if cu == "—" and ta == "Y":
            ta_color = VIOLET
        add_text(s, tx + Inches(6.3), ry, Inches(1.5), row_h,
                 cu, size=14, color=cu_color, bold=True, font="Calibri",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + Inches(7.9), ry, Inches(1.5), row_h,
                 ta, size=14, color=ta_color, bold=True, font="Calibri",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + Inches(9.6), ry, Inches(2.5), row_h,
                 note, size=10, color=MUTED, italic=True, font="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 8)
    set_notes(s,
        "Position vs the competitor everyone knows. 'Cursor is the closest "
        "analogue. So we measured. Fifteen capabilities — composer, "
        "codebase index, at-mentions, MCP servers, model switcher, usage "
        "console, magic-link, privacy mode — that's Cursor parity, "
        "shipped. Below the fold, ten things Cursor structurally cannot "
        "ship without redirecting their product: voice cockpit, "
        "hash-chained receipts, Solana anchor, seven life-domain packs, "
        "the meeet.world economy. We play the broader board. Cursor plays "
        "the deeper board. The gap matters: regulated buyers, life-ops "
        "buyers, on-chain buyers — none of those addressable to Cursor.'")
    return s


def slide_9():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Differentiators",
              "Five things Cursor cannot ship.",
              title_size=34)

    add_text(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
             "Structural — not roadmap items.",
             size=13, color=MUTED, italic=True, font="Calibri")

    items = [
        ("01", "Voice-driven Composer",
         "Cmd+Shift+Space → voice prompt → multi-file diff → approve → receipt. "
         "Same loop, no keyboard.",
         "Cursor is keyboard-bound. We're hands-free.", CYAN),
        ("02", "Receipt-anchored audit + Solana proof",
         "Every consequential action emits a signed receipt; daily Merkle root "
         "anchored on Solana memo program; public verifier with zero auth.",
         "Cursor has no answer to 'prove it'.", VIOLET),
        ("03", "SOC2 + GDPR compliance bundle",
         "compliance_export router, audit log, BAA-ready posture. Regulated "
         "industries can deploy on-prem with the policy bundle pre-built.",
         "Cursor enterprise exists; regulated tier doesn't.", INDIGO),
        ("04", "7 domain packs — not just code",
         "Wealth, health, business, entrepreneur, research, civic, creative. "
         "Each pack ships its vocabulary, actions, connectors day-one.",
         "Cursor is an IDE. There is no Wealth pack for IDEs.", LIME),
        ("05", "$MEEET economy through meeet.world",
         "Pay in USD or $MEEET at 10% discount. Marketplace 70/30 split, "
         "on-chain payouts. CAC leverage in the crypto-native segment.",
         "Cursor has Stripe. We have rails.", AMBER),
    ]

    y = Inches(2.05); row_h = Inches(0.88); gap = Inches(0.08)
    for (num, title, body, kicker, color) in items:
        add_rect(s, Inches(0.6), y, Inches(12.2), row_h,
                 fill=PANEL, line=LINE, rounded=True)
        add_rect(s, Inches(0.8), y + Inches(0.2),
                 Inches(0.55), Inches(0.55), fill=color, rounded=True)
        add_text(s, Inches(0.8), y + Inches(0.2),
                 Inches(0.55), Inches(0.55), num,
                 size=14, color=BG, bold=True, font="Consolas",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.55), y + Inches(0.1),
                 Inches(7), Inches(0.4),
                 title, size=16, color=WHITE, bold=True, font="Calibri")
        add_text(s, Inches(1.55), y + Inches(0.45),
                 Inches(7.5), Inches(0.5),
                 body, size=11, color=BODY, font="Calibri", line_spacing=1.3)
        add_text(s, Inches(9.3), y + Inches(0.28),
                 Inches(3.4), Inches(0.5),
                 kicker, size=11, color=color, italic=True, bold=True,
                 font="Calibri", line_spacing=1.3,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + gap

    footer(s, 9)
    set_notes(s,
        "Slide 9 is the moat slide. Walk it fast — these are five sentences. "
        "'One: voice-driven Composer. Two: receipt-anchored audit and the "
        "Solana proof — this is the answer to the regulator and the partner "
        "asking did the agent really do that. Three: SOC2 + GDPR. Four: "
        "seven life-domain packs — Cursor cannot ship a wealth pack without "
        "leaving the IDE and disowning their distribution. Five: $MEEET. "
        "Each of these is structural, not a roadmap item. Cursor can't ship "
        "any of them in twelve months without disowning their core.' Land "
        "the line, then go to pricing.")
    return s


def slide_10():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Pricing", "Three tiers. Two currencies. One ledger.",
              title_size=30)

    px = Inches(0.6); pw = Inches(2.85); py = Inches(1.8); ph = Inches(4.0); gap = Inches(0.20)

    tiers = [
        ("FREE", "$0", "/mo", "Civic pack always free",
         ["50 requests / month",
          "Basic models only",
          "Watermarked receipts",
          "Civic pack — free for all"],
         DIM),
        ("PRO", "$20", "/mo or 200 $MEEET",
         "Solo operators",
         ["1,000 requests / month",
          "All providers · BYO key",
          "Clean receipts + Solana anchor",
          "Marketplace earn-side enabled"],
         CYAN),
        ("BUSINESS", "$40", "/seat or 400 $MEEET",
         "Funds · ops · regulated",
         ["5,000 req / seat (soft cap)",
          "SOC2-style audit log",
          "Unlimited workspaces",
          "On-prem option · $1k+/seat"],
         VIOLET),
    ]

    for i, (name, price, period, sub, feats, color) in enumerate(tiers):
        x = px + i * (pw + gap)
        is_mid = (i == 1)
        card_fill = PANEL_2 if is_mid else PANEL
        card_line = color if is_mid else LINE
        add_rect(s, x, py, pw, ph, fill=card_fill, line=card_line,
                 line_w=2.0 if is_mid else 1.0, rounded=True)
        add_text(s, x + Inches(0.3), py + Inches(0.3),
                 pw - Inches(0.6), Inches(0.4),
                 name, size=14, color=color, bold=True, font="Consolas")
        add_text(s, x + Inches(0.3), py + Inches(0.75),
                 pw - Inches(0.6), Inches(0.95),
                 price, size=52, color=WHITE, bold=True, font="Calibri")
        add_text(s, x + Inches(0.3), py + Inches(1.65),
                 pw - Inches(0.6), Inches(0.32),
                 period, size=10, color=MUTED, font="Consolas")
        add_text(s, x + Inches(0.3), py + Inches(2.00),
                 pw - Inches(0.6), Inches(0.30),
                 sub, size=11, color=BODY, italic=True, font="Calibri")
        add_rect(s, x + Inches(0.3), py + Inches(2.40),
                 pw - Inches(0.6), Inches(0.02), fill=LINE)
        fy = py + Inches(2.55)
        for f in feats:
            add_rect(s, x + Inches(0.3), fy + Inches(0.10),
                     Inches(0.08), Inches(0.08), fill=color, rounded=True)
            add_text(s, x + Inches(0.5), fy,
                     pw - Inches(0.8), Inches(0.32),
                     f, size=10, color=BODY, font="Calibri",
                     line_spacing=1.3)
            fy += Inches(0.34)

    rx = Inches(0.6) + 3 * (pw + gap)
    rw = SLIDE_W - rx - Inches(0.6)
    add_rect(s, rx, py, rw, ph,
             fill=PANEL, line=AMBER, rounded=True)
    add_text(s, rx + Inches(0.3), py + Inches(0.3),
             rw - Inches(0.6), Inches(0.4),
             "RECEIPT FLOW", size=12, color=AMBER, bold=True,
             font="Consolas")
    flow = [
        "1.  Agent action",
        "2.  Signed receipt",
        "3.  Hash chain",
        "4.  Daily Merkle root",
        "5.  Solana memo anchor",
        "6.  Public verifier",
    ]
    fy = py + Inches(0.85)
    for line in flow:
        add_text(s, rx + Inches(0.3), fy, rw - Inches(0.6), Inches(0.35),
                 line, size=12, color=BODY, font="Consolas")
        fy += Inches(0.42)

    footer(s, 10)
    set_notes(s,
        "Business model. Four streams. 'Self-serve subscriptions — $20/mo "
        "PRO, $40/seat BUSINESS. Same shape Cursor uses, same shape the "
        "market underwrites without question. Alternative payment in "
        "$MEEET at a 10% discount — drives ecosystem usage, reduces CAC "
        "in the crypto-native segment, never mandatory. On-prem licensing "
        "at $1k-plus per seat per month for regulated buyers — highest "
        "LTV stream we can charge enterprise-SaaS rates on. And the "
        "marketplace, opening v10.5 — we take 30% on third-party skills. "
        "Every line on the receipt flow on the right is shipped today.'")
    return s


def slide_11():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "B2B angle",
              "Regulated industries Cursor cannot enter.",
              title_size=30)

    rx = Inches(8.7); ry = Inches(1.9); rw = Inches(4.1); rh = Inches(4.8)
    add_rect(s, rx, ry, rw, rh, fill=PANEL, line=VIOLET, line_w=2.0,
             rounded=True)
    add_text(s, rx, ry + Inches(0.4), rw, Inches(0.5),
             "COMPLIANCE POSTURE", size=11, color=VIOLET, bold=True,
             font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(1.0), rw, Inches(1.0),
             "SOC2", size=80, color=WHITE, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(2.1), rw, Inches(0.45),
             "Type II", size=22, color=CYAN, bold=True, italic=True,
             font="Calibri", align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(2.7), rw, Inches(0.35),
             "+  GDPR exporter", size=14, color=BODY, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(3.1), rw, Inches(0.35),
             "+  BAA-ready posture", size=14, color=BODY, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(3.5), rw, Inches(0.35),
             "+  audit log + receipts", size=14, color=BODY, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, rx, ry + Inches(4.0), rw, Inches(0.35),
             "W257  ·  shipped", size=10, color=DIM, bold=True,
             font="Consolas", align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.6), Inches(1.85), Inches(8), Inches(0.4),
             "TARGET BUYERS", size=11, color=CYAN, bold=True,
             font="Consolas")

    sectors = [
        ("Hedge & PE funds",
         "On-prem TARS for analyst desks. Algotrade pack + audit-grade receipts "
         "satisfy investor reporting and compliance teams.",
         AMBER),
        ("Healthcare ops",
         "Local-first means PHI never leaves the box. BAA posture + GDPR "
         "exporter unlock provider IT review.",
         LIME),
        ("Legal & accounting",
         "Domain pack covers research, citation graph, receipt-anchored "
         "evidence trail. Bar-association friendly defaults.",
         CYAN),
        ("Government & civic tech",
         "Civic pack free for all tiers. Local-first deployment, "
         "FedRAMP-style roadmap, public verifier as a transparency primitive.",
         INDIGO),
    ]
    y = Inches(2.4)
    for (title, body, color) in sectors:
        add_rect(s, Inches(0.6), y, Inches(7.85), Inches(1.0),
                 fill=PANEL, line=LINE, rounded=True)
        add_rect(s, Inches(0.78), y + Inches(0.32),
                 Inches(0.16), Inches(0.36), fill=color)
        add_text(s, Inches(1.05), y + Inches(0.12),
                 Inches(6.6), Inches(0.4),
                 title, size=15, color=WHITE, bold=True, font="Calibri")
        add_text(s, Inches(1.05), y + Inches(0.45),
                 Inches(6.7), Inches(0.55),
                 body, size=10, color=BODY, font="Calibri", line_spacing=1.3)
        y += Inches(1.10)

    footer(s, 11)
    set_notes(s,
        "Enterprise revenue thesis. 'This is the slide that moves the room "
        "if your investor cares about ARR. Cursor is locked out of "
        "regulated buyers by design — VS Code fork in a cloud-default "
        "posture is a non-starter for a hedge fund's IT review. We ship "
        "SOC2 Type II posture today — the readiness doc, the GDPR "
        "exporter, the BAA-ready stance, the audit log, and the "
        "compliance bundle export. On-prem licensing at $1k-plus per "
        "seat-month is the highest-LTV stream. Four buyer segments — "
        "funds, healthcare, legal, government — each one a Cursor-blind "
        "spot. Three are already in the pipeline.'")
    return s


def slide_12():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "The cockpit", "Cinematic. Ambient. One shortcut away.",
              title_size=28)

    px = Inches(0.6); py = Inches(1.95); pw = Inches(6.05); ph = Inches(2.50); gap = Inches(0.18)

    panels = [
        ("Monolith — idle/listening",
         "Cyan-violet light strip · breathing · 200ms response",
         CYAN),
        ("Audit Explorer",
         "Receipt #1284 · hash chain · Solana memo TX hash",
         AMBER),
        ("Composer panel — voice diff",
         "Multi-file diff · domain-pack aware · approve hunks",
         VIOLET),
        ("USAGE tab",
         "Live token meter · soft cap warning · $MEEET balance",
         LIME),
    ]

    positions = [
        (px, py),
        (px + pw + gap, py),
        (px, py + ph + gap),
        (px + pw + gap, py + ph + gap),
    ]

    for (caption_title, caption_body, color), (x, y) in zip(panels, positions):
        add_rect(s, x, y, pw, ph, fill=PANEL, line=LINE, line_w=1.0,
                 rounded=True)
        ix = x + Inches(0.15); iy = y + Inches(0.15)
        iw = pw - Inches(0.3); ih = ph - Inches(0.95)
        add_rect(s, ix, iy, iw, ih, fill=BG, line=color, line_w=1.0,
                 rounded=True)
        for i, dx in enumerate([0.0, 0.06, 0.12]):
            add_rect(s, ix + Inches(0.2 + dx), iy + Inches(0.2),
                     Inches(0.025), Inches(0.20),
                     fill=color)
        add_text(s, ix, iy, iw, ih,
                 "[ screenshot placeholder ]",
                 size=13, color=DIM, italic=True, font="Consolas",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), y + ph - Inches(0.75),
                 pw - Inches(0.5), Inches(0.35),
                 caption_title, size=13, color=WHITE, bold=True,
                 font="Calibri")
        add_text(s, x + Inches(0.25), y + ph - Inches(0.40),
                 pw - Inches(0.5), Inches(0.30),
                 caption_body, size=10, color=color, italic=True,
                 font="Consolas")

    footer(s, 12)
    set_notes(s,
        "Show the surfaces. Replace these four placeholders with real "
        "screenshots from TARS.app the morning of the demo — they will "
        "kill the live deck if they're stale. Capture order: one, the "
        "monolith in listening state with the cyan-violet strip visible; "
        "two, Audit Explorer with at least one receipt visible plus the "
        "Solana TX hash badge; three, the Composer panel with an actual "
        "diff in view; four, the Usage tab showing the live token meter "
        "and the $MEEET balance. Same crop, same dark theme. Don't show "
        "any dev-tools overlay. Capture at 2x then scale for retina.")
    return s


def slide_13():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Code transparency", "Open development. Main branch only.",
              title_size=32)

    sw = Inches(2.95); sh = Inches(1.55); sy = Inches(2.0); sgap = Inches(0.20)
    stat_callout(s, Inches(0.6), sy, sw, sh,
                 "53", "commits W203→W269", CYAN)
    stat_callout(s, Inches(0.6) + (sw + sgap), sy, sw, sh,
                 "~33k", "lines of code", VIOLET)
    stat_callout(s, Inches(0.6) + 2 * (sw + sgap), sy, sw, sh,
                 "50", "FastAPI routers", INDIGO)
    stat_callout(s, Inches(0.6) + 3 * (sw + sgap), sy, sw, sh,
                 "100%", "main branch", LIME)

    add_rect(s, Inches(0.6), Inches(4.0), Inches(12.2), Inches(2.2),
             fill=PANEL, line=LINE, rounded=True)
    add_text(s, Inches(1.0), Inches(4.25), Inches(11.4), Inches(0.45),
             "EVERY ACTION", size=12, color=AMBER, bold=True,
             font="Consolas")
    add_text(s, Inches(1.0), Inches(4.65), Inches(11.4), Inches(0.6),
             "emits a hash-chained receipt. Anchored on Solana daily.",
             size=22, color=WHITE, bold=True, italic=True, font="Calibri")
    add_text(s, Inches(1.0), Inches(5.35), Inches(11.4), Inches(0.85),
             "github.com/alienram/jarvis · MIT-licensed backend · "
             "open development · audit trail in main branch · "
             "every receipt verifiable by a stranger with no auth.",
             size=13, color=BODY, font="Calibri", line_spacing=1.4)

    footer(s, 13)
    set_notes(s,
        "Open development as proof of velocity and trust. 'Fifty-three "
        "commits between W203 and W269. Roughly thirty-three thousand "
        "lines of code. Fifty FastAPI routers. Everything in main. No "
        "private repo, no hidden branch — investors and partners can "
        "read every line. The receipt primitive is not a marketing "
        "claim; it is wired into the orchestrator and visible at "
        "github dot com slash alienram slash jarvis. The MIT license on "
        "the backend means an on-prem buyer can audit the substrate and "
        "fork it if we ever stop being trustworthy. That is the offer.'")
    return s


def slide_14():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Roadmap", "v10.1 → v11.0  ·  Always-on substrate.",
              title_size=30)

    versions = [
        ("v10.0", "Q2-26", "GA",         "Voice cockpit · receipts · 7 packs", CYAN),
        ("v10.1", "Q3-26", "Polish",     "Mobile companion · perf · 5 langs",  INDIGO),
        ("v10.2", "Q3-26", "Plugins",    "Marketplace v1 · 70/30 payouts",          VIOLET),
        ("v10.5", "Q4-26", "Workshop",   "B2B cohort dashboard · on-prem GTM",      AMBER),
        ("v10.8", "Q1-27", "Federation", "T2T agent handshake · multi-tenant",      LIME),
        ("v11.0", "Q2-27", "Agentic OS", "Always-on substrate · daemons · OS hooks", ROSE),
    ]

    line_y = Inches(4.0)
    add_rect(s, Inches(1.4), line_y, Inches(10.5), Inches(0.04), fill=LINE)

    n = len(versions)
    total_w = Inches(10.5)
    spacing = total_w / (n - 1)

    for i, (ver, qtr, theme, body, color) in enumerate(versions):
        cx = Inches(1.4) + i * spacing
        marker_r = Inches(0.18)
        add_rect(s, cx - marker_r / 2, line_y - marker_r / 2 + Inches(0.02),
                 marker_r, marker_r, fill=color, rounded=True)
        if i % 2 == 0:
            cy = Inches(1.85)
        else:
            cy = Inches(4.30)
        cw = Inches(1.75); ch = Inches(2.05)
        cx_card = cx - cw / 2

        add_rect(s, cx_card, cy, cw, ch, fill=PANEL, line=color, line_w=1.0,
                 rounded=True)
        add_text(s, cx_card + Inches(0.15), cy + Inches(0.15),
                 cw - Inches(0.3), Inches(0.4),
                 ver, size=18, color=WHITE, bold=True, font="Calibri")
        add_text(s, cx_card + Inches(0.15), cy + Inches(0.55),
                 cw - Inches(0.3), Inches(0.3),
                 qtr, size=10, color=color, bold=True, font="Consolas")
        add_text(s, cx_card + Inches(0.15), cy + Inches(0.85),
                 cw - Inches(0.3), Inches(0.4),
                 theme, size=12, color=WHITE, bold=True, italic=True,
                 font="Calibri")
        add_text(s, cx_card + Inches(0.15), cy + Inches(1.25),
                 cw - Inches(0.3), Inches(1.0),
                 body, size=9, color=BODY, font="Calibri", line_spacing=1.3)

    footer(s, 14)
    set_notes(s,
        "Where we're going. Read the headline themes, not the dates. "
        "'v10.0 ships tomorrow. v10.1 is polish — mobile companion, "
        "performance, full five-language i18n. v10.2 opens the "
        "marketplace, third parties publishing skills with 70/30 "
        "payouts on Solana. v10.5 is the B2B push — workshop cohort "
        "dashboard, on-prem deployments at fund-sized contracts. v10.8 "
        "lights up federation — agent-to-agent handshake across "
        "TARS instances. v11.0 is the inflection: always-on agentic "
        "substrate. Not a chat surface — the OS layer for agentic work. "
        "This is the eighteen-month bridge.'")
    return s


def slide_15():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Brother sync", "meeet.world integration — 2 weeks to live billing.",
              title_size=28)

    add_text(s, Inches(0.6), Inches(1.95), Inches(7), Inches(0.4),
             "ENDPOINTS BROTHER IS WIRING", size=11, color=CYAN, bold=True,
             font="Consolas")

    endpoints = [
        "POST  /api/auth/meeet/exchange",
        "POST  /api/billing/usage_event",
        "GET   /api/billing/balance",
        "POST  /api/billing/charge",
        "GET   /api/entitlements/check",
        "POST  /api/marketplace/payout",
        "POST  /api/identity/oauth/callback",
        "POST  /api/anchor/merkle_batch",
    ]
    y = Inches(2.40)
    for ep in endpoints:
        add_rect(s, Inches(0.6), y, Inches(7), Inches(0.42),
                 fill=PANEL, line=LINE, rounded=True)
        add_text(s, Inches(0.85), y, Inches(0.45), Inches(0.42),
                 ">", size=12, color=CYAN, bold=True, font="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.25), y, Inches(5.6), Inches(0.42),
                 ep, size=11, color=BODY, font="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.0), y, Inches(0.9), Inches(0.42),
                 "wired", size=9, color=LIME, bold=True, font="Consolas",
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.50)

    rx = Inches(8.1); rw = Inches(4.7)
    add_rect(s, rx, Inches(1.95), rw, Inches(2.0), fill=PANEL,
             line=AMBER, line_w=2.0, rounded=True)
    add_text(s, rx, Inches(2.10), rw, Inches(0.4),
             "ETA TO LIVE BILLING", size=11, color=AMBER, bold=True,
             font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, rx, Inches(2.50), rw, Inches(0.95),
             "2 weeks", size=54, color=WHITE, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, rx, Inches(3.45), rw, Inches(0.45),
             "TARS-side ready · brother's 4 cloud endpoints in test",
             size=11, color=BODY, italic=True, font="Calibri",
             align=PP_ALIGN.CENTER)

    add_rect(s, rx, Inches(4.15), rw, Inches(2.35), fill=PANEL,
             line=VIOLET, line_w=2.0, rounded=True)
    add_text(s, rx, Inches(4.30), rw, Inches(0.4),
             "WHY THIS MATTERS", size=11, color=VIOLET, bold=True,
             font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.3), Inches(4.75), rw - Inches(0.6), Inches(1.7),
             "meeet.world becomes the universal identity, balance, and "
             "compliance substrate. One brother-owned domain authenticates "
             "every TARS instance — and any third-party app on the same "
             "rails. The marketplace, the on-chain economy, and the "
             "regulator handshake all flow through one place.",
             size=11, color=BODY, font="Calibri", line_spacing=1.4)

    footer(s, 15)
    set_notes(s,
        "Partner story. 'My brother runs meeet.world. It's not a vendor "
        "relationship — it's a coordinated build. He owns identity, "
        "billing, the $MEEET token economy, and the marketplace ledger. "
        "TARS owns the cockpit and the receipt primitive. Eight endpoints "
        "shape the contract — all eight wired on the TARS side, four "
        "still in test on his side. Live billing flips on in two weeks. "
        "The bigger story: meeet.world becomes the universal substrate "
        "for any agent-action product that needs identity, balance, "
        "compliance telemetry. We're the first tenant. Not the last.'")
    return s


def slide_16():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Traction", "v10.0 GA tomorrow.  Pipeline below.",
              title_size=32)

    sx = Inches(0.6); sw = Inches(3.95); sy = Inches(2.0); sh = Inches(2.3); gap = Inches(0.21)

    add_rect(s, sx, sy, sw, sh, fill=PANEL, line=CYAN, line_w=2.0, rounded=True)
    add_text(s, sx, sy + Inches(0.35), sw, Inches(0.4),
             "v10.0 GA", size=11, color=CYAN, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)
    add_text(s, sx, sy + Inches(0.85), sw, Inches(1.0),
             "SHIPPED", size=46, color=WHITE, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, sx, sy + Inches(1.65), sw, Inches(0.5),
             "rc.1 cut · GA gate: tomorrow", size=12, color=BODY,
             italic=True, font="Calibri", align=PP_ALIGN.CENTER)

    add_rect(s, sx + sw + gap, sy, sw, sh, fill=PANEL, line=VIOLET, line_w=2.0,
             rounded=True)
    add_text(s, sx + sw + gap, sy + Inches(0.35), sw, Inches(0.4),
             "WAITLIST", size=11, color=VIOLET, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)
    add_text(s, sx + sw + gap, sy + Inches(0.85), sw, Inches(1.0),
             "200+", size=64, color=WHITE, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, sx + sw + gap, sy + Inches(1.85), sw, Inches(0.3),
             "verified email signups", size=11, color=MUTED,
             italic=True, font="Calibri", align=PP_ALIGN.CENTER)

    add_rect(s, sx + 2 * (sw + gap), sy, sw, sh, fill=PANEL, line=AMBER,
             line_w=2.0, rounded=True)
    add_text(s, sx + 2 * (sw + gap), sy + Inches(0.35), sw, Inches(0.4),
             "ON-PREM LEADS", size=11, color=AMBER, bold=True,
             font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, sx + 2 * (sw + gap), sy + Inches(0.85), sw, Inches(1.0),
             "3", size=64, color=WHITE, bold=True, font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, sx + 2 * (sw + gap), sy + Inches(1.85), sw, Inches(0.3),
             "warm · fund · healthcare · gov", size=11, color=MUTED,
             italic=True, font="Calibri", align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.6), Inches(5.0), Inches(12.2), Inches(1.5),
             fill=PANEL, line=LINE, rounded=True)
    add_text(s, Inches(0.85), Inches(5.15), Inches(11.5), Inches(0.4),
             "MOMENTUM", size=11, color=LIME, bold=True, font="Consolas")
    add_text(s, Inches(0.85), Inches(5.55), Inches(11.5), Inches(0.45),
             "267 commits  ·  10 months  ·  one founder + AI orchestra  ·  "
             "v10.0 GA on schedule",
             size=18, color=WHITE, bold=True, font="Calibri")
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.5), Inches(0.4),
             "Velocity is the moat investors can underwrite without a market call.",
             size=12, color=MUTED, italic=True, font="Calibri")

    footer(s, 16)
    set_notes(s,
        "Show momentum, not vanity metrics. Replace the 200 waitlist "
        "number with the live count the morning of the pitch — pull from "
        "the production dashboard. Replace 3 on-prem leads with whatever "
        "is qualified that day. Speak the story: 'v10.0 GA ships "
        "tomorrow. Two hundred verified waitlist signups, three warm "
        "on-prem leads — a fund, a healthcare org, a civic-tech buyer. "
        "What matters more than absolute numbers is velocity: 267 "
        "commits in ten months, by one founder running an AI orchestra. "
        "Same discipline scales the next 18 months.'")
    return s


def slide_17():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Team", "One founder.  One brother.  One AI orchestra.",
              title_size=30)

    cards = [
        ("ALIEN",         "Founder · CEO",
         "Voice of the product. Vision, design, lane orchestration. "
         "Shipped 267 commits in 10 months with the AI orchestra. "
         "Operator first, builder always.",
         CYAN),
        ("BROTHER",       "meeet.world · Substrate",
         "Owns identity, billing, marketplace ledger, $MEEET economy. "
         "Eight endpoints away from live billing. Independent runway, "
         "shared protocol surface.",
         VIOLET),
        ("AI ORCHESTRA",  "Claude · Cursor lanes",
         "Treated as engineering capacity, not a tool. Claude lane "
         "leads design and core ops. Cursor lane handles "
         "implementation. Managed like an offshore team.",
         INDIGO),
    ]

    px = Inches(0.6); pw = Inches(4.04); py = Inches(2.0); ph = Inches(3.5); gap = Inches(0.18)
    for i, (name, role, body, color) in enumerate(cards):
        x = px + i * (pw + gap)
        add_rect(s, x, py, pw, ph, fill=PANEL, line=LINE, rounded=True)
        add_rect(s, x + Inches(0.3), py + Inches(0.3),
                 Inches(0.9), Inches(0.9), fill=color, rounded=True)
        add_text(s, x + Inches(0.3), py + Inches(0.3),
                 Inches(0.9), Inches(0.9), name[0],
                 size=40, color=BG, bold=True, font="Calibri",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.4), py + Inches(0.35),
                 pw - Inches(1.6), Inches(0.5),
                 name, size=20, color=WHITE, bold=True, font="Calibri")
        add_text(s, x + Inches(1.4), py + Inches(0.85),
                 pw - Inches(1.6), Inches(0.35),
                 role, size=11, color=color, bold=True, italic=True,
                 font="Consolas")
        add_text(s, x + Inches(0.3), py + Inches(1.5),
                 pw - Inches(0.6), ph - Inches(1.7),
                 body, size=12, color=BODY, font="Calibri", line_spacing=1.4)

    add_rect(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(0.9),
             fill=PANEL, line=AMBER, rounded=True)
    add_text(s, Inches(0.85), Inches(5.75), Inches(11.5), Inches(0.4),
             "CAPITAL EFFICIENCY", size=11, color=AMBER, bold=True,
             font="Consolas")
    add_text(s, Inches(0.85), Inches(6.10), Inches(11.5), Inches(0.5),
             "53 commits in one session · v10 GA on a tiny burn · "
             "the structure is the story",
             size=15, color=WHITE, bold=True, italic=True, font="Calibri")

    footer(s, 17)
    set_notes(s,
        "Capital efficiency story. 'This is a tiny team by design. One "
        "founder. One brother running the adjacent substrate. An AI "
        "orchestra — Claude and Cursor — that we manage like an offshore "
        "team. The structure is the story. We built v10 GA, with audit "
        "trails and on-prem deployment, at the headcount of a typical "
        "seed-stage team that shipped half as much. Series A funds the "
        "next layer — distribution, sales, brother's runway, and three "
        "first hires. Not engineering — the engineering is already a "
        "compounding asset.'")
    return s


def slide_18():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "The ask", "Three concrete things.",
              title_size=42)

    asks = [
        ("01", "Design partners",
         "Three companies for the v10.1 cohort. Funds, healthcare ops, "
         "or civic-tech preferred — anywhere the receipt primitive moves a "
         "compliance conversation.",
         "Outcome: shaped roadmap + case studies + references.",
         CYAN),
        ("02", "On-prem pilots",
         "Two pilots at $1k+/seat/month, six-month minimum. We bring the "
         "compliance bundle, the workshop pack, and white-glove "
         "deployment.",
         "Outcome: $12-25k MRR per pilot · highest-LTV stream.",
         VIOLET),
        ("03", "$MEEET pre-orders",
         "Pre-order $MEEET at the launch peg. Drives ecosystem activity "
         "and gives early users a discounted runway. Opt-in, not "
         "mandatory.",
         "Outcome: ecosystem velocity · CAC reduction · brother's runway.",
         AMBER),
    ]
    y = Inches(2.0); row_h = Inches(1.35); gap = Inches(0.15)
    for (num, title, body, outcome, color) in asks:
        add_rect(s, Inches(0.6), y, Inches(12.2), row_h,
                 fill=PANEL, line=color, line_w=1.5, rounded=True)
        add_rect(s, Inches(0.85), y + Inches(0.3),
                 Inches(0.85), Inches(0.85), fill=color, rounded=True)
        add_text(s, Inches(0.85), y + Inches(0.3),
                 Inches(0.85), Inches(0.85), num,
                 size=28, color=BG, bold=True, font="Consolas",
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.9), y + Inches(0.18),
                 Inches(6.5), Inches(0.45),
                 title, size=22, color=WHITE, bold=True, font="Calibri")
        add_text(s, Inches(1.9), y + Inches(0.65),
                 Inches(6.5), Inches(0.65),
                 body, size=11, color=BODY, font="Calibri", line_spacing=1.35)
        add_text(s, Inches(8.6), y + Inches(0.35),
                 Inches(4.0), Inches(0.65),
                 outcome, size=11, color=color, italic=True, bold=True,
                 font="Calibri", line_spacing=1.35,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += row_h + gap

    footer(s, 18)
    set_notes(s,
        "Close. Three specific things. 'We're not raising a round in this "
        "meeting — we're asking for the things that compound. One: three "
        "design partners for the v10.1 cohort. If you know a fund, a "
        "healthcare org, or a civic-tech buyer where receipts move a "
        "compliance conversation, introduce us. Two: two on-prem pilots "
        "at thousand-dollar-per-seat. Six months, white glove. Three: "
        "$MEEET pre-orders at the launch peg — opt-in, never mandatory. "
        "Optional but it funds my brother's runway and seeds the "
        "ecosystem. Those are the three asks. Series A conversation "
        "happens after T+30 traction.'")
    return s


def slide_19():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)
    title_bar(s, "Q&A prep · presenter notes",
              "Five anticipated.  One line each.",
              title_size=28)

    qas = [
        ("What if Cursor enters this market?",
         "Cursor is a VS Code fork. To enter, they redirect product or "
         "extend IDE to non-devs — neither is a 12-month bet.", VIOLET),
        ("Why local-first if inference is cloud anyway?",
         "Data plane separation. Receipts, memory, index never leave the "
         "box. Inference call is the only network egress, with BYO key.", CYAN),
        ("Is $MEEET a security?",
         "Utility token for compute on the meeet.world relayer. "
         "Day-one users can ignore the token. On-prem buyers don't touch it.", AMBER),
        ("How big is the moat without inference cost advantage?",
         "Moat is the cockpit, the receipt primitive, the marketplace, the "
         "per-user style — not inference price. Cheap inference helps us "
         "asymmetrically.", LIME),
        ("Why hasn't an incumbent built this?",
         "Anthropic/OpenAI ship chat surfaces. Microsoft Copilot is wired into "
         "Office. None ship voice-first + local + receipts + 7 packs together.", INDIGO),
    ]

    y = Inches(1.85)
    h = Inches(0.95)
    for (q, a, color) in qas:
        add_rect(s, Inches(0.6), y, Inches(12.2), h,
                 fill=PANEL, line=LINE, rounded=True)
        add_rect(s, Inches(0.78), y + Inches(0.18),
                 Inches(0.12), h - Inches(0.36), fill=color)
        add_text(s, Inches(1.05), y + Inches(0.10),
                 Inches(11.4), Inches(0.4),
                 "Q.  " + q, size=13, color=WHITE, bold=True,
                 font="Calibri")
        add_text(s, Inches(1.05), y + Inches(0.45),
                 Inches(11.4), Inches(0.5),
                 "A.  " + a, size=10, color=BODY, italic=True,
                 font="Calibri", line_spacing=1.35)
        y += h + Inches(0.05)

    add_text(s, Inches(0.6), SLIDE_H - Inches(0.85),
             Inches(12.2), Inches(0.35),
             "Hide this slide before live delivery — presenter notes only.",
             size=10, color=AMBER, italic=True, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)

    footer(s, 19)
    set_notes(s,
        "PRESENTER-ONLY slide. Hide before live delivery — File > Slide "
        "Properties > Hidden. If you forget to hide it, no disaster: the "
        "answers are honest and concise. Anchor every answer to one of the "
        "five lines on the slide rather than improvising. Practiced "
        "version: Cursor objection answered first, $MEEET regulatory "
        "answered cleanest. Don't speak more than ten seconds per answer "
        "in live Q&A — let the room ask follow-ups instead.")
    return s


def slide_20():
    s = prs.slides.add_slide(blank)
    background(s)
    vignette_corners(s)

    draw_monolith(s, SLIDE_W / 2, Inches(3.0), Inches(3.5), w_ratio=0.30)
    draw_ring(s, SLIDE_W / 2, Inches(3.0), Inches(1.8),
              color=RGBColor(0x22, 0x24, 0x4c), thickness_pt=0.5)
    draw_ring(s, SLIDE_W / 2, Inches(3.0), Inches(2.4),
              color=RGBColor(0x16, 0x18, 0x35), thickness_pt=0.5)

    add_text(s, Inches(0.6), Inches(4.85), Inches(12.2), Inches(0.8),
             "Thank you.", size=54, color=WHITE, bold=True, italic=True,
             font="Calibri", align=PP_ALIGN.CENTER)

    add_runs(s, Inches(0.6), Inches(5.65), Inches(12.2), Inches(0.5),
             [("Build the cockpit ", 16, BODY, False, False, "Calibri"),
              ("for everything not code.", 16, VIOLET, True, True, "Calibri")],
             align=PP_ALIGN.CENTER)

    cy = Inches(6.4)
    add_text(s, Inches(0.6), cy, Inches(12.2), Inches(0.4),
             "alienram@icloud.com   ·   meeet.world   ·   github.com/alienram/jarvis",
             size=13, color=CYAN, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), cy + Inches(0.35), Inches(12.2), Inches(0.3),
             "TARS  ·  v10.0.0-rc.1  ·  ships tomorrow",
             size=10, color=DIM, font="Consolas", align=PP_ALIGN.CENTER)

    set_notes(s,
        "Final word. Pause. Make eye contact across the room. 'Cursor "
        "did this for code. We're doing this for everything else. "
        "v10 ships tomorrow. We'd love to find out together what an AI "
        "cockpit looks like when it owns the whole desk, not just the "
        "editor. The contact line is on the screen — email, the brother's "
        "domain, the open repository. Thank you.' Then stop talking. "
        "Don't fill the silence. Let the room respond. The first "
        "question is the one to listen to most carefully — it tells you "
        "what the room actually heard.")
    return s


# Build all slides
for builder in [slide_1, slide_2, slide_3, slide_4, slide_5,
                slide_6, slide_7, slide_8, slide_9, slide_10,
                slide_11, slide_12, slide_13, slide_14, slide_15,
                slide_16, slide_17, slide_18, slide_19, slide_20]:
    builder()

import sys
out_path = sys.argv[1] if len(sys.argv) > 1 else "TARS_v10.0_PRESENTATION.pptx"
prs.save(out_path)
print(f"Wrote {out_path} - {len(prs.slides)} slides")
